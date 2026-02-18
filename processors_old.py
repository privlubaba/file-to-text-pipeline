"""
File Processors Module - Enhanced Version
Handles extraction of text from different file formats with support for:
- Complex tables with structure preservation
- Embedded images with OCR
- Quality assessment metrics
- Enhanced metadata extraction
"""

import io
from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional, Tuple
import logging
import base64
import json
from pathlib import Path

# PDF processing
import pdfplumber
from pypdf import PdfReader
from pdf2image import convert_from_bytes
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter

# Office documents
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import openpyxl

# Utilities
import subprocess
import tempfile
import os
import re
from collections import defaultdict

# Language detection (optional)
try:
    from langdetect import detect, LangDetectException
    LANGDETECT_AVAILABLE = True
except ImportError:
    LANGDETECT_AVAILABLE = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# Arabic text utilities
class ArabicTextNormalizer:
    """Utilities for normalizing Arabic text"""

    # Arabic diacritical marks (Harakat/Tashkeel)
    ARABIC_DIACRITICS = [
        '\u064B',  # Fathatan
        '\u064C',  # Dammatan
        '\u064D',  # Kasratan
        '\u064E',  # Fatha
        '\u064F',  # Damma
        '\u0650',  # Kasra
        '\u0651',  # Shadda
        '\u0652',  # Sukun
        '\u0653',  # Maddah
        '\u0654',  # Hamza above
        '\u0655',  # Hamza below
        '\u0656',  # Subscript alef
        '\u0657',  # Inverted damma
        '\u0658',  # Mark noon ghunna
        '\u0670',  # Superscript alef
    ]

    @staticmethod
    def remove_diacritics(text: str) -> str:
        """Remove Arabic diacritics (harakat) from text"""
        for diacritic in ArabicTextNormalizer.ARABIC_DIACRITICS:
            text = text.replace(diacritic, '')
        return text

    @staticmethod
    def normalize_arabic(text: str) -> str:
        """Normalize Arabic text for better readability"""
        # Remove diacritics
        text = ArabicTextNormalizer.remove_diacritics(text)

        # Normalize Arabic letters
        # Replace different forms of Alef
        text = re.sub('[إأآا]', 'ا', text)
        # Replace Teh Marbuta with Heh
        text = text.replace('ة', 'ه')
        # Replace different forms of Yeh
        text = re.sub('[ىي]', 'ي', text)

        return text

    @staticmethod
    def clean_spacing(text: str) -> str:
        """Clean up spacing issues in Arabic text"""
        # Remove multiple spaces
        text = re.sub(r'\s+', ' ', text)
        # Remove spaces before punctuation
        text = re.sub(r'\s+([.,!?،؛])', r'\1', text)
        return text.strip()


class QualityMetrics:
    """Quality assessment metrics for extracted content"""
    def __init__(self):
        self.total_text_length = 0
        self.total_tables = 0
        self.total_images = 0
        self.average_confidence = 0.0
        self.language_consistency = 1.0
        self.structure_score = 0.0
        self.completeness_score = 0.0

    def calculate_overall_score(self) -> float:
        """Calculate overall quality score (0-100)"""
        scores = []

        # Text extraction score (has meaningful content)
        if self.total_text_length > 100:
            scores.append(min(100, (self.total_text_length / 1000) * 50))

        # Structure score (tables and formatting preserved)
        scores.append(self.structure_score)

        # Confidence score (OCR accuracy)
        scores.append(self.average_confidence)

        # Completeness score
        scores.append(self.completeness_score)

        return sum(scores) / len(scores) if scores else 0.0

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return {
            "overall_score": round(self.calculate_overall_score(), 2),
            "total_text_length": self.total_text_length,
            "total_tables_extracted": self.total_tables,
            "total_images_processed": self.total_images,
            "average_ocr_confidence": round(self.average_confidence, 2),
            "language_consistency": round(self.language_consistency, 2),
            "structure_preservation_score": round(self.structure_score, 2),
            "completeness_score": round(self.completeness_score, 2),
            "quality_rating": self._get_quality_rating()
        }

    def _get_quality_rating(self) -> str:
        """Get human-readable quality rating"""
        score = self.calculate_overall_score()
        if score >= 90:
            return "Excellent"
        elif score >= 75:
            return "Good"
        elif score >= 60:
            return "Fair"
        elif score >= 40:
            return "Poor"
        else:
            return "Very Poor"


class TextExtractionResult:
    """Container for text extraction results with enhanced features"""
    def __init__(self, file_id: str, filename: str, file_type: str):
        self.file_id = file_id
        self.filename = filename
        self.file_type = file_type
        self.pages: List[Dict[str, Any]] = []
        self.total_pages = 0
        self.extraction_method = ""
        self.metadata: Dict[str, Any] = {}
        self.quality_metrics = QualityMetrics()
        self.tables: List[Dict[str, Any]] = []
        self.images: List[Dict[str, Any]] = []

    def add_page(self, page_number: int, text: str, method: str, **kwargs):
        """Add a page of extracted text"""
        # Detect if text contains Arabic
        has_arabic = bool(re.search(r'[\u0600-\u06FF]', text))

        page_data = {
            "page_number": page_number,
            "raw_text": text,
            "char_count": len(text),
            "extraction_method": method,
            **kwargs
        }

        # Add normalized Arabic text if Arabic is detected
        if has_arabic:
            normalized_text = ArabicTextNormalizer.remove_diacritics(text)
            cleaned_text = ArabicTextNormalizer.clean_spacing(normalized_text)
            page_data["normalized_text"] = cleaned_text
            page_data["has_diacritics"] = text != normalized_text
            page_data["language"] = "arabic"

        self.pages.append(page_data)
        self.total_pages = len(self.pages)

        # Update quality metrics
        self.quality_metrics.total_text_length += len(text)

    def add_table(self, page_number: int, table_data: List[List[str]], table_index: int = 0):
        """Add extracted table with structure"""
        self.tables.append({
            "page_number": page_number,
            "table_index": table_index,
            "rows": len(table_data),
            "columns": len(table_data[0]) if table_data else 0,
            "data": table_data
        })
        self.quality_metrics.total_tables += 1
        self.quality_metrics.structure_score = min(100, self.quality_metrics.total_tables * 20)

    def add_image(self, page_number: int, image_index: int, image_info: Dict[str, Any]):
        """Add extracted image information"""
        self.images.append({
            "page_number": page_number,
            "image_index": image_index,
            **image_info
        })
        self.quality_metrics.total_images += 1

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization"""
        return {
            "file_id": self.file_id,
            "filename": self.filename,
            "file_type": self.file_type,
            "total_pages": self.total_pages,
            "extraction_method": self.extraction_method,
            "metadata": self.metadata,
            "quality_assessment": self.quality_metrics.to_dict(),
            "pages": self.pages,
            "tables": self.tables if self.tables else [],
            "images": self.images if self.images else []
        }


class FileProcessor(ABC):
    """Abstract base class for file processors"""
    
    @abstractmethod
    def can_process(self, mime_type: str, extension: str) -> bool:
        """Check if this processor can handle the file type"""
        pass
    
    @abstractmethod
    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        """Extract text from the file"""
        pass


class PDFProcessor(FileProcessor):
    """Process PDF files with OCR support for scanned documents"""
    
    def can_process(self, mime_type: str, extension: str) -> bool:
        return mime_type == 'application/pdf' or extension.lower() == '.pdf'
    
    def _get_available_languages(self) -> List[str]:
        """Get list of available Tesseract languages"""
        try:
            result = subprocess.run(
                ['tesseract', '--list-langs'],
                capture_output=True,
                text=True
            )
            langs = result.stdout.split('\n')[1:]  # Skip header
            available = [lang.strip() for lang in langs if lang.strip() and lang.strip() != 'osd']
            logger.info(f"Available Tesseract languages: {', '.join(available)}")
            return available
        except Exception as e:
            logger.warning(f"Could not get Tesseract languages: {e}")
            return ['eng']  # Default to English
    
    def _detect_language_from_image(self, file_bytes: bytes) -> str:
        """Detect language from scanned PDF by testing OCR with different languages"""
        try:
            # Convert first page to image
            images = convert_from_bytes(file_bytes, dpi=200, first_page=1, last_page=1)
            if not images:
                return 'eng'
            
            image = images[0]
            # Preprocess
            image = self._preprocess_image(image)
            
            # Test with multiple languages and pick the one with best confidence
            test_langs = ['ara', 'eng', 'fra', 'spa']  # Common languages
            available_langs = self._get_available_languages()
            test_langs = [lang for lang in test_langs if lang in available_langs]
            
            best_lang = 'eng'
            best_score = 0
            
            logger.info(f"Testing OCR with languages: {', '.join(test_langs)}")
            
            for lang in test_langs:
                try:
                    # Quick OCR test
                    result = pytesseract.image_to_data(
                        image,
                        lang=lang,
                        output_type=pytesseract.Output.DICT,
                        config='--psm 6'
                    )
                    
                    # Calculate score based on confidence and character count
                    confidences = [int(conf) for conf in result['conf'] if conf != '-1']
                    texts = [text for text in result['text'] if text.strip()]
                    
                    if confidences and texts:
                        avg_conf = sum(confidences) / len(confidences)
                        char_count = sum(len(text) for text in texts)
                        # Score = confidence * log(char_count)
                        import math
                        score = avg_conf * math.log(max(char_count, 1))
                        
                        logger.info(f"  {lang}: confidence={avg_conf:.1f}, chars={char_count}, score={score:.1f}")
                        
                        if score > best_score:
                            best_score = score
                            best_lang = lang
                except Exception as e:
                    logger.warning(f"  {lang}: failed - {e}")
                    continue
            
            logger.info(f"Selected language: {best_lang} (score: {best_score:.1f})")
            return best_lang
            
        except Exception as e:
            logger.warning(f"Image language detection failed: {e}")
            return 'eng'
    
    def _detect_language(self, file_bytes: bytes, is_scanned: bool = False) -> str:
        """Detect document language automatically"""
        
        # For scanned PDFs, use image-based detection
        if is_scanned:
            logger.info("Scanned PDF detected, using image-based language detection")
            return self._detect_language_from_image(file_bytes)
        
        # For text-based PDFs, use text-based detection
        if not LANGDETECT_AVAILABLE:
            logger.info("langdetect not available, using English")
            return 'eng'
        
        try:
            # Try to extract some text first
            reader = PdfReader(io.BytesIO(file_bytes))
            sample_text = ""
            
            # Get text from first few pages
            for i in range(min(3, len(reader.pages))):
                sample_text += reader.pages[i].extract_text()
                if len(sample_text) > 100:
                    break
            
            # If we have text, detect language
            if sample_text.strip():
                detected = detect(sample_text)
                
                # Map langdetect codes to Tesseract language codes
                lang_map = {
                    'ar': 'ara',      # Arabic
                    'en': 'eng',      # English
                    'fr': 'fra',      # French
                    'es': 'spa',      # Spanish
                    'de': 'deu',      # German
                    'zh-cn': 'chi_sim',  # Chinese Simplified
                    'zh-tw': 'chi_tra',  # Chinese Traditional
                    'ja': 'jpn',      # Japanese
                    'ko': 'kor',      # Korean
                    'ru': 'rus',      # Russian
                    'hi': 'hin',      # Hindi
                    'ur': 'urd',      # Urdu
                }
                
                tesseract_lang = lang_map.get(detected, 'eng')
                
                # Check if language is available
                available_langs = self._get_available_languages()
                if tesseract_lang not in available_langs:
                    logger.warning(f"Language '{tesseract_lang}' detected but not installed. Using English. Install with: brew install tesseract-lang (macOS) or sudo apt-get install tesseract-ocr-{detected[:3]} (Linux)")
                    return 'eng'
                
                logger.info(f"Detected language from text: {detected} -> Tesseract: {tesseract_lang}")
                return tesseract_lang
            
        except Exception as e:
            logger.warning(f"Text-based language detection failed: {e}, using English")
        
        # Default to English
        return 'eng'
    
    def _is_scanned_pdf(self, reader: PdfReader) -> bool:
        """Detect if PDF is scanned (image-based)"""
        try:
            # Check first page for text
            first_page_text = reader.pages[0].extract_text().strip()
            return len(first_page_text) < 50  # If very little text, likely scanned
        except:
            return True
    
    def _extract_tables_from_page(self, page) -> List[List[List[str]]]:
        """Extract structured tables from a PDF page"""
        tables = []
        try:
            extracted_tables = page.extract_tables()
            if extracted_tables:
                for table in extracted_tables:
                    # Clean and structure table data
                    cleaned_table = []
                    for row in table:
                        if row:
                            cleaned_row = [str(cell).strip() if cell else "" for cell in row]
                            cleaned_table.append(cleaned_row)
                    if cleaned_table:
                        tables.append(cleaned_table)
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")
        return tables

    def _extract_images_from_page(self, page, page_num: int) -> List[Dict[str, Any]]:
        """Extract images from a PDF page"""
        images_info = []
        try:
            if hasattr(page, 'images'):
                for img_idx, img in enumerate(page.images):
                    image_info = {
                        "width": img.get('width', 0),
                        "height": img.get('height', 0),
                        "x0": img.get('x0', 0),
                        "y0": img.get('y0', 0),
                        "type": "embedded_image"
                    }
                    images_info.append(image_info)
        except Exception as e:
            logger.warning(f"Image extraction failed for page {page_num}: {e}")
        return images_info

    def _extract_text_based(self, file_bytes: bytes, result: TextExtractionResult) -> List[str]:
        """Extract text from text-based PDF using pdfplumber with enhanced features"""
        texts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                # Extract text
                text = page.extract_text() or ""
                texts.append(text)

                # Extract tables
                tables = self._extract_tables_from_page(page)
                for table_idx, table in enumerate(tables):
                    result.add_table(page_num, table, table_idx)

                # Extract images
                images = self._extract_images_from_page(page, page_num)
                for img_idx, img_info in enumerate(images):
                    result.add_image(page_num, img_idx, img_info)

        return texts
    
    def _extract_with_ocr(self, file_bytes: bytes, language: str = 'eng', result: TextExtractionResult = None) -> List[str]:
        """Extract text from scanned PDF using OCR with quality metrics"""
        texts = []
        confidences = []
        try:
            # Convert PDF to images with higher DPI for better quality
            images = convert_from_bytes(file_bytes, dpi=400)

            logger.info(f"Using OCR language: {language}")

            for i, image in enumerate(images):
                # Preprocess image for better OCR
                processed_image = self._preprocess_image(image)

                # Apply OCR with detected language and get detailed data
                custom_config = r'--oem 3 --psm 6'  # LSTM OCR Engine, assume uniform text block

                # Get text with confidence scores
                ocr_data = pytesseract.image_to_data(
                    processed_image,
                    lang=language,
                    config=custom_config,
                    output_type=pytesseract.Output.DICT
                )

                # Extract text
                text = pytesseract.image_to_string(
                    processed_image,
                    lang=language,
                    config=custom_config
                )

                # Calculate confidence for this page
                page_confidences = [int(conf) for conf in ocr_data['conf'] if conf != '-1']
                page_avg_confidence = sum(page_confidences) / len(page_confidences) if page_confidences else 0
                confidences.append(page_avg_confidence)

                # Clean up the text
                text = text.strip()
                texts.append(text)

                # Store image info
                if result:
                    result.add_image(i + 1, 0, {
                        "width": image.width,
                        "height": image.height,
                        "type": "scanned_page",
                        "ocr_confidence": round(page_avg_confidence, 2),
                        "text_length": len(text)
                    })

                logger.info(f"OCR processed page {i+1}: {len(text)} characters extracted (confidence: {page_avg_confidence:.1f}%)")

            # Update quality metrics
            if result and confidences:
                result.quality_metrics.average_confidence = sum(confidences) / len(confidences)
                result.quality_metrics.completeness_score = min(100, result.quality_metrics.total_text_length / 100)

            logger.info(f"OCR processed {len(images)} pages total")
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
            raise

        return texts
    
    def _preprocess_image(self, image):
        """Preprocess image for better OCR results"""
        from PIL import ImageEnhance, ImageFilter
        
        # Convert to grayscale
        image = image.convert('L')
        
        # Increase contrast
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.0)
        
        # Sharpen
        image = image.filter(ImageFilter.SHARPEN)
        
        return image
    
    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        """Extract text from PDF (with automatic OCR detection)"""
        result = TextExtractionResult(file_id, filename, "pdf")
        
        try:
            # Try to detect if it's scanned
            reader = PdfReader(io.BytesIO(file_bytes))
            is_scanned = self._is_scanned_pdf(reader)
            
            # Detect language (pass is_scanned for better detection)
            detected_lang = self._detect_language(file_bytes, is_scanned=is_scanned)
            
            # Extract metadata
            result.metadata = {
                "page_count": len(reader.pages),
                "is_scanned": is_scanned,
                "detected_language": detected_lang,
                "pdf_metadata": {
                    "title": reader.metadata.title if reader.metadata else None,
                    "author": reader.metadata.author if reader.metadata else None,
                }
            }
            
            # Choose extraction method
            if is_scanned:
                logger.info(f"Detected scanned PDF, using OCR ({detected_lang}) for {filename}")
                texts = self._extract_with_ocr(file_bytes, detected_lang, result)
                method = f"tesseract_ocr_{detected_lang}"
            else:
                logger.info(f"Detected text-based PDF, using direct extraction for {filename}")
                texts = self._extract_text_based(file_bytes, result)
                method = "pdfplumber"
                # For text-based PDFs, set high confidence and completeness
                result.quality_metrics.average_confidence = 100.0
                result.quality_metrics.completeness_score = min(100, result.quality_metrics.total_text_length / 100)
            
            result.extraction_method = method
            
            # Store per-page text
            total_chars = 0
            for i, text in enumerate(texts, start=1):
                result.add_page(i, text, method)
                total_chars += len(text.strip())
            
            # Check if extraction was successful
            if total_chars == 0:
                result.metadata["warning"] = "No text could be extracted from this document"
                result.metadata["possible_reasons"] = [
                    "Document contains only images with no readable text",
                    "Document is password protected",
                    "OCR could not recognize any text (poor image quality)",
                    "Document is completely blank",
                    f"Required language pack '{detected_lang}' may not be installed"
                ]
                logger.warning(f"No text extracted from {filename} (0 characters)")
            elif total_chars < 50 and is_scanned:
                result.metadata["warning"] = "Very little text extracted - OCR may have struggled with image quality"
                result.metadata["suggestions"] = [
                    "Try scanning the original document at higher resolution (300+ DPI)",
                    "Ensure the document is not skewed or rotated",
                    "Check if the image is clear and high contrast",
                    f"If document is in {detected_lang}, ensure language pack is installed: brew install tesseract-lang"
                ]
            
            logger.info(f"Successfully extracted {len(texts)} pages from PDF: {filename} ({total_chars} total characters, language: {detected_lang})")
            
        except Exception as e:
            logger.error(f"Failed to process PDF {filename}: {e}")
            raise
        
        return result


class DOCXProcessor(FileProcessor):
    """Process DOCX (Word) files with enhanced table and image extraction"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return (mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                              'application/msword'] or
                extension.lower() in ['.docx', '.doc'])

    def _extract_tables_from_docx(self, doc) -> List[List[List[str]]]:
        """Extract structured tables from DOCX"""
        tables_data = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                tables_data.append(table_data)
        return tables_data

    def _extract_images_from_docx(self, doc) -> List[Dict[str, Any]]:
        """Extract image information from DOCX"""
        images_info = []
        try:
            # Access the document's relationships to find images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    images_info.append({
                        "type": "embedded_image",
                        "format": rel.target_ref.split('.')[-1],
                        "relationship_id": rel.rId
                    })
        except Exception as e:
            logger.warning(f"Image extraction from DOCX failed: {e}")
        return images_info

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        """Extract text from DOCX file with enhanced features"""
        result = TextExtractionResult(file_id, filename, "docx")
        result.extraction_method = "python-docx-enhanced"

        try:
            # Use python-docx for extraction
            doc = DocxDocument(io.BytesIO(file_bytes))

            # Extract all paragraphs
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)

            # Extract tables with structure
            tables = self._extract_tables_from_docx(doc)
            for table_idx, table in enumerate(tables):
                result.add_table(1, table, table_idx)
                # Add table as text representation too
                table_text = "\n[TABLE]\n"
                for row in table:
                    table_text += " | ".join(row) + "\n"
                full_text.append(table_text)

            # Extract images
            images = self._extract_images_from_docx(doc)
            for img_idx, img_info in enumerate(images):
                result.add_image(1, img_idx, img_info)

            # Join all text
            complete_text = '\n'.join(full_text)

            # Store as single "page" (Word docs don't have fixed pages)
            result.add_page(1, complete_text, "python-docx-enhanced")

            # Update quality metrics
            result.quality_metrics.average_confidence = 100.0  # Native format, high confidence
            result.quality_metrics.completeness_score = min(100, len(complete_text) / 100)

            result.metadata = {
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "image_count": len(images),
                "total_characters": len(complete_text),
                "has_complex_content": len(tables) > 0 or len(images) > 0
            }

            logger.info(f"Successfully extracted text from DOCX: {filename} ({len(tables)} tables, {len(images)} images)")

        except Exception as e:
            logger.error(f"Failed to process DOCX {filename}: {e}")
            # Fallback: try pandoc if available
            try:
                result = self._extract_with_pandoc(file_bytes, filename, file_id)
            except:
                raise e

        return result
    
    def _extract_with_pandoc(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        """Fallback extraction using pandoc"""
        result = TextExtractionResult(file_id, filename, "docx")
        result.extraction_method = "pandoc"
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
            temp_file.write(file_bytes)
            temp_file.flush()
            
            try:
                # Run pandoc
                output = subprocess.check_output(
                    ['pandoc', '--track-changes=all', temp_file.name, '-t', 'plain'],
                    stderr=subprocess.STDOUT
                )
                text = output.decode('utf-8')
                result.add_page(1, text, "pandoc")
            finally:
                os.unlink(temp_file.name)
        
        return result


class PPTXProcessor(FileProcessor):
    """Process PPTX (PowerPoint) files with enhanced extraction"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return (mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or
                extension.lower() == '.pptx')

    def _extract_tables_from_slide(self, slide) -> List[List[List[str]]]:
        """Extract tables from a slide"""
        tables = []
        try:
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        tables.append(table_data)
        except Exception as e:
            logger.warning(f"Table extraction from slide failed: {e}")
        return tables

    def _count_images_in_slide(self, slide) -> int:
        """Count images in a slide"""
        image_count = 0
        try:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    image_count += 1
        except:
            pass
        return image_count

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        """Extract text from PPTX file (slide by slide) with enhanced features"""
        result = TextExtractionResult(file_id, filename, "pptx")
        result.extraction_method = "python-pptx-enhanced"

        try:
            prs = Presentation(io.BytesIO(file_bytes))
            total_images = 0

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                # Extract tables
                tables = self._extract_tables_from_slide(slide)
                for table_idx, table in enumerate(tables):
                    result.add_table(slide_num, table, table_idx)
                    # Add table representation to text
                    table_text = "\n[TABLE]\n"
                    for row in table:
                        table_text += " | ".join(row) + "\n"
                    slide_text.append(table_text)

                # Count images
                img_count = self._count_images_in_slide(slide)
                total_images += img_count
                if img_count > 0:
                    for img_idx in range(img_count):
                        result.add_image(slide_num, img_idx, {
                            "type": "slide_image",
                            "format": "unknown"
                        })

                # Extract speaker notes if present
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_text.append(f"\n[Speaker Notes]: {notes_text}")

                # Combine all text for this slide
                combined_text = '\n'.join(slide_text)
                result.add_page(slide_num, combined_text, "python-pptx-enhanced",
                              slide_number=slide_num,
                              images_count=img_count,
                              tables_count=len(tables))

            # Update quality metrics
            result.quality_metrics.average_confidence = 100.0  # Native format
            result.quality_metrics.completeness_score = min(100, result.quality_metrics.total_text_length / 100)

            result.metadata = {
                "slide_count": len(prs.slides),
                "total_slides": len(prs.slides),
                "total_images": total_images,
                "has_complex_content": result.quality_metrics.total_tables > 0 or total_images > 0
            }

            logger.info(f"Successfully extracted text from PPTX: {filename} ({len(prs.slides)} slides, {result.quality_metrics.total_tables} tables, {total_images} images)")

        except Exception as e:
            logger.error(f"Failed to process PPTX {filename}: {e}")
            raise

        return result


class XLSXProcessor(FileProcessor):
    """Process XLSX (Excel) files"""
    
    def can_process(self, mime_type: str, extension: str) -> bool:
        return (mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                              'application/vnd.ms-excel'] or
                extension.lower() in ['.xlsx', '.xls'])
    
    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        """Extract text from XLSX file (sheet by sheet)"""
        result = TextExtractionResult(file_id, filename, "xlsx")
        result.extraction_method = "pandas"
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
            
            for sheet_num, sheet_name in enumerate(excel_file.sheet_names, start=1):
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert dataframe to text representation
                sheet_text = f"Sheet: {sheet_name}\n\n"
                sheet_text += df.to_string(index=False)
                
                result.add_page(
                    sheet_num, 
                    sheet_text, 
                    "pandas",
                    sheet_name=sheet_name,
                    rows=len(df),
                    columns=len(df.columns)
                )
            
            result.metadata = {
                "sheet_count": len(excel_file.sheet_names),
                "sheet_names": excel_file.sheet_names
            }
            
            logger.info(f"Successfully extracted text from XLSX: {filename} ({len(excel_file.sheet_names)} sheets)")
            
        except Exception as e:
            logger.error(f"Failed to process XLSX {filename}: {e}")
            raise
        
        return result


class TXTProcessor(FileProcessor):
    """Process plain text files"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return mime_type == 'text/plain' or extension.lower() == '.txt'

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> 'TextExtractionResult':
        """Extract text from TXT file"""
        try:
            # Try UTF-8 first
            text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1
            try:
                text = file_bytes.decode('latin-1')
            except:
                text = file_bytes.decode('utf-8', errors='ignore')

        # Create page
        page = PageText(
            page_number=1,
            raw_text=text,
            char_count=len(text),
            extraction_method="plain_text"
        )

        # Detect language
        language = self._detect_language(text)
        page.language = language

        # Normalize Arabic if detected
        if language == 'ar':
            page.has_diacritics = self._has_arabic_diacritics(text)
            if page.has_diacritics:
                page.normalized_text = ArabicTextNormalizer.remove_diacritics(text)

        return TextExtractionResult(
            file_id=file_id,
            filename=filename,
            file_type="txt",
            total_pages=1,
            extraction_method="plain_text",
            pages=[page],
            metadata={"encoding": "utf-8", "file_size": len(file_bytes)}
        )


class HTMLProcessor(FileProcessor):
    """Process HTML files"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return mime_type == 'text/html' or extension.lower() in ['.html', '.htm']

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> 'TextExtractionResult':
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            HAS_BS4 = True
        except ImportError:
            HAS_BS4 = False

        try:
            html_content = file_bytes.decode('utf-8')
        except:
            html_content = file_bytes.decode('utf-8', errors='ignore')

        if HAS_BS4:
            # Use BeautifulSoup for better parsing
            soup = BeautifulSoup(html_content, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator='\n')
        else:
            # Simple regex fallback
            text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL)
            text = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL)
            text = re.sub(r'<[^>]+>', '', text)

        # Clean up whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text.strip()

        page = PageText(
            page_number=1,
            raw_text=text,
            char_count=len(text),
            extraction_method="html_parser"
        )

        language = self._detect_language(text)
        page.language = language

        if language == 'ar':
            page.has_diacritics = self._has_arabic_diacritics(text)
            if page.has_diacritics:
                page.normalized_text = ArabicTextNormalizer.remove_diacritics(text)

        return TextExtractionResult(
            file_id=file_id,
            filename=filename,
            file_type="html",
            total_pages=1,
            extraction_method="html_parser",
            pages=[page],
            metadata={"file_size": len(file_bytes)}
        )


class MarkdownProcessor(FileProcessor):
    """Process Markdown files"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return mime_type == 'text/markdown' or extension.lower() in ['.md', '.markdown']

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> 'TextExtractionResult':
        """Extract text from Markdown file"""
        try:
            text = file_bytes.decode('utf-8')
        except:
            text = file_bytes.decode('utf-8', errors='ignore')

        # Remove markdown formatting (basic)
        # Remove headers
        text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
        # Remove bold/italic
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        # Remove links but keep text
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
        # Remove images
        text = re.sub(r'!\[([^\]]*)\]\([^\)]+\)', '', text)
        # Remove code blocks
        text = re.sub(r'```[^`]*```', '', text, flags=re.DOTALL)
        text = re.sub(r'`([^`]+)`', r'\1', text)

        page = PageText(
            page_number=1,
            raw_text=text,
            char_count=len(text),
            extraction_method="markdown_parser"
        )

        language = self._detect_language(text)
        page.language = language

        if language == 'ar':
            page.has_diacritics = self._has_arabic_diacritics(text)
            if page.has_diacritics:
                page.normalized_text = ArabicTextNormalizer.remove_diacritics(text)

        return TextExtractionResult(
            file_id=file_id,
            filename=filename,
            file_type="markdown",
            total_pages=1,
            extraction_method="markdown_parser",
            pages=[page],
            metadata={"file_size": len(file_bytes)}
        )


class RTFProcessor(FileProcessor):
    """Process RTF (Rich Text Format) files"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return mime_type in ['text/rtf', 'application/rtf'] or extension.lower() == '.rtf'

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> 'TextExtractionResult':
        """Extract text from RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            HAS_STRIPRTF = True
        except ImportError:
            HAS_STRIPRTF = False

        try:
            rtf_content = file_bytes.decode('utf-8')
        except:
            rtf_content = file_bytes.decode('latin-1', errors='ignore')

        if HAS_STRIPRTF:
            text = rtf_to_text(rtf_content)
        else:
            # Basic RTF parsing fallback
            text = re.sub(r'\\[a-z]+\d*\s?', '', rtf_content)
            text = re.sub(r'[{}]', '', text)
            text = text.strip()

        page = PageText(
            page_number=1,
            raw_text=text,
            char_count=len(text),
            extraction_method="rtf_parser"
        )

        language = self._detect_language(text)
        page.language = language

        if language == 'ar':
            page.has_diacritics = self._has_arabic_diacritics(text)
            if page.has_diacritics:
                page.normalized_text = ArabicTextNormalizer.remove_diacritics(text)

        return TextExtractionResult(
            file_id=file_id,
            filename=filename,
            file_type="rtf",
            total_pages=1,
            extraction_method="rtf_parser",
            pages=[page],
            metadata={"file_size": len(file_bytes)}
        )


class ODTProcessor(FileProcessor):
    """Process ODT (OpenDocument Text) files"""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return mime_type == 'application/vnd.oasis.opendocument.text' or extension.lower() == '.odt'

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> 'TextExtractionResult':
        """Extract text from ODT file"""
        import zipfile
        from xml.etree import ElementTree as ET

        try:
            # ODT files are ZIP archives
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as odt_file:
                # Read content.xml
                content_xml = odt_file.read('content.xml')
                root = ET.fromstring(content_xml)

                # Extract text from all text nodes
                text_parts = []
                for elem in root.iter():
                    if elem.text:
                        text_parts.append(elem.text)
                    if elem.tail:
                        text_parts.append(elem.tail)

                text = '\n'.join(text_parts)
                text = re.sub(r'\n\s*\n', '\n\n', text).strip()

        except Exception as e:
            logger.error(f"Error extracting ODT: {e}")
            text = f"Error extracting ODT content: {str(e)}"

        page = PageText(
            page_number=1,
            raw_text=text,
            char_count=len(text),
            extraction_method="odt_parser"
        )

        language = self._detect_language(text)
        page.language = language

        if language == 'ar':
            page.has_diacritics = self._has_arabic_diacritics(text)
            if page.has_diacritics:
                page.normalized_text = ArabicTextNormalizer.remove_diacritics(text)

        return TextExtractionResult(
            file_id=file_id,
            filename=filename,
            file_type="odt",
            total_pages=1,
            extraction_method="odt_parser",
            pages=[page],
            metadata={"file_size": len(file_bytes)}
        )


class ProcessorFactory:
    """Factory for creating appropriate file processors"""
    
    def __init__(self):
        self.processors = [
            PDFProcessor(),
            DOCXProcessor(),
            PPTXProcessor(),
            XLSXProcessor(),
            TXTProcessor(),
            HTMLProcessor(),
            MarkdownProcessor(),
            RTFProcessor(),
            ODTProcessor(),
        ]
    
    def get_processor(self, mime_type: str, extension: str) -> Optional[FileProcessor]:
        """Get the appropriate processor for a file type"""
        for processor in self.processors:
            if processor.can_process(mime_type, extension):
                return processor
        return None