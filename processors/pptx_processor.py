"""
Enhanced PPTX Processor.
Recursive shape group traversal, SmartArt text, chart titles, speaker notes.
"""

import io
import logging
from typing import List, Dict, Any

from processors.base import FileProcessor, TextExtractionResult

logger = logging.getLogger(__name__)


class PPTXProcessor(FileProcessor):
    """Process PPTX files with enhanced shape, SmartArt, and chart extraction."""

    def can_process(self, mime_type: str, extension: str) -> bool:
        return (mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                or extension.lower() == '.pptx')

    def extract_text(self, file_bytes: bytes, filename: str, file_id: str) -> TextExtractionResult:
        from pptx import Presentation
        from pptx.util import Inches

        result = TextExtractionResult(file_id, filename, "pptx")
        result.extraction_method = "python-pptx-enhanced"

        try:
            prs = Presentation(io.BytesIO(file_bytes))
            total_images = 0

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []

                # Extract text from all shapes recursively
                for shape in slide.shapes:
                    shape_texts = self._extract_text_from_shape(shape)
                    slide_text.extend(shape_texts)

                # Extract tables
                tables = self._extract_tables_from_slide(slide)
                for table_idx, table in enumerate(tables):
                    result.add_table(slide_num, table, table_idx)
                    table_text = "\n[TABLE]\n" + "\n".join(" | ".join(row) for row in table)
                    slide_text.append(table_text)

                # Count and record images
                img_count = self._extract_images_from_slide(slide, slide_num, result)
                total_images += img_count

                # Extract speaker notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_text.append(f"\n[SPEAKER NOTES]: {notes_text}")

                combined_text = '\n'.join(slide_text)
                result.add_page(slide_num, combined_text, "python-pptx-enhanced",
                                slide_number=slide_num,
                                images_count=img_count,
                                tables_count=len(tables))

            result.quality_metrics.average_confidence = 100.0
            result.quality_metrics.completeness_score = min(
                100, result.quality_metrics.total_text_length / 100
            )

            result.metadata = {
                "slide_count": len(prs.slides),
                "total_images": total_images,
                "total_tables": result.quality_metrics.total_tables
            }

            logger.info(
                f"Extracted text from PPTX: {filename} "
                f"({len(prs.slides)} slides, {result.quality_metrics.total_tables} tables, "
                f"{total_images} images)"
            )

        except Exception as e:
            logger.error(f"Failed to process PPTX {filename}: {e}")
            raise

        return result

    def _extract_text_from_shape(self, shape, depth: int = 0) -> List[str]:
        """
        Recursively extract text from shapes, including grouped shapes,
        SmartArt, and chart titles.
        """
        texts = []
        max_depth = 10  # Prevent infinite recursion

        if depth > max_depth:
            return texts

        # Handle group shapes — recurse into children
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child_shape in shape.shapes:
                    texts.extend(self._extract_text_from_shape(child_shape, depth + 1))
            except Exception:
                pass

        # Regular text from text frame
        if hasattr(shape, "text_frame"):
            try:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        texts.append(para_text)
            except Exception:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)

        # SmartArt text — parse underlying DrawingML XML
        if hasattr(shape, '_element'):
            smartart_texts = self._extract_smartart_text(shape._element)
            for st in smartart_texts:
                if st not in texts:  # Avoid duplicates
                    texts.append(st)

        # Chart titles and labels
        if shape.has_chart:
            chart_texts = self._extract_chart_text(shape.chart)
            texts.extend(chart_texts)

        return texts

    def _extract_smartart_text(self, element) -> List[str]:
        """Extract text from SmartArt/DrawingML XML elements."""
        texts = []
        try:
            ns = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
            }
            # Look for text elements in DrawingML
            for t_elem in element.findall('.//a:t', ns):
                if t_elem.text and t_elem.text.strip():
                    texts.append(t_elem.text.strip())
        except Exception:
            pass
        return texts

    def _extract_chart_text(self, chart) -> List[str]:
        """Extract chart title, axis labels, and series names."""
        texts = []
        try:
            # Chart title
            if chart.chart_title and chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text.strip()
                if title_text:
                    texts.append(f"[CHART TITLE]: {title_text}")

            # Series names
            for series in chart.series:
                try:
                    if hasattr(series, 'tx') and series.tx:
                        # Access series name from the XML
                        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
                        for v_elem in series.tx._element.findall('.//c:v', ns):
                            if v_elem.text:
                                texts.append(f"[CHART SERIES]: {v_elem.text}")
                except Exception:
                    pass

            # Category labels
            try:
                plot = chart.plots[0]
                if hasattr(plot, 'categories') and plot.categories:
                    cats = list(plot.categories)
                    if cats:
                        texts.append(f"[CHART CATEGORIES]: {', '.join(str(c) for c in cats)}")
            except Exception:
                pass

        except Exception as e:
            logger.debug(f"Chart text extraction failed: {e}")
        return texts

    def _extract_tables_from_slide(self, slide) -> List[List[List[str]]]:
        """Extract tables from a slide."""
        tables = []
        try:
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        tables.append(table_data)
        except Exception as e:
            logger.debug(f"Table extraction from slide failed: {e}")
        return tables

    def _extract_images_from_slide(self, slide, slide_num: int,
                                    result: TextExtractionResult) -> int:
        """Extract image metadata from slide. Returns count of images found."""
        image_count = 0
        try:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image_count += 1
                    img_info = {
                        "type": "slide_image",
                        "width": shape.width,
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top
                    }
                    # Try to get image format
                    try:
                        if hasattr(shape, 'image') and shape.image:
                            img_info["format"] = shape.image.content_type
                    except Exception:
                        img_info["format"] = "unknown"

                    result.add_image(slide_num, image_count - 1, img_info)
        except Exception:
            pass
        return image_count
